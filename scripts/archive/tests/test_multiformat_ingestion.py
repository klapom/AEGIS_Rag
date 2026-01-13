#!/usr/bin/env python3
"""Multi-Format Ingestion Test Script.

This script tests document ingestion across all supported formats
and collects detailed metrics for analysis.

Usage:
    python scripts/test_multiformat_ingestion.py
    python scripts/test_multiformat_ingestion.py --format pdf
    python scripts/test_multiformat_ingestion.py --all
"""

import argparse
import json
import os
import sys
import tempfile
import time
from datetime import datetime
from pathlib import Path
from typing import Any

import requests

# Configuration
API_BASE_URL = os.getenv("API_BASE_URL", "http://localhost:8000")
TEST_USERNAME = os.getenv("TEST_USERNAME", "admin")
TEST_PASSWORD = os.getenv("TEST_PASSWORD", "admin123")

# Test content
TEST_CONTENT = {
    "title": "AegisRAG Multi-Format Test Document",
    "sections": [
        {
            "heading": "Introduction to Machine Learning",
            "content": """Machine learning is a subset of artificial intelligence that enables
            systems to learn and improve from experience. Key concepts include supervised learning,
            unsupervised learning, and reinforcement learning. Popular frameworks include
            TensorFlow, PyTorch, and scikit-learn.""",
        },
        {
            "heading": "Natural Language Processing",
            "content": """NLP focuses on the interaction between computers and human language.
            Important techniques include tokenization, named entity recognition (NER),
            sentiment analysis, and transformer models like BERT and GPT. SpaCy and
            Hugging Face are widely used libraries.""",
        },
        {
            "heading": "Knowledge Graphs",
            "content": """Knowledge graphs represent information as entities and relationships.
            Neo4j is a popular graph database. Entity extraction identifies people, organizations,
            and concepts. Relationship extraction finds connections between entities.
            AegisRAG uses LightRAG for graph-based retrieval.""",
        },
    ],
    "table_data": [
        ["Framework", "Type", "Language"],
        ["TensorFlow", "Deep Learning", "Python"],
        ["PyTorch", "Deep Learning", "Python"],
        ["SpaCy", "NLP", "Python"],
        ["Neo4j", "Graph Database", "Java"],
    ],
}


def get_auth_token() -> str:
    """Get authentication token."""
    print(f"[AUTH] Authenticating as {TEST_USERNAME}...")
    response = requests.post(
        f"{API_BASE_URL}/api/v1/auth/login",
        json={"username": TEST_USERNAME, "password": TEST_PASSWORD},
        timeout=30,
    )
    if response.status_code != 200:
        print(f"[AUTH] FAILED: {response.text}")
        sys.exit(1)
    token = response.json()["access_token"]
    print(f"[AUTH] Success!")
    return token


def create_csv_file(temp_dir: Path) -> Path:
    """Create sample CSV file."""
    csv_path = temp_dir / "test_document.csv"
    lines = [",".join(TEST_CONTENT["table_data"][0])]
    for row in TEST_CONTENT["table_data"][1:]:
        lines.append(",".join(row))
    lines.append("")
    lines.append("# Document Content")
    for section in TEST_CONTENT["sections"]:
        lines.append(f"## {section['heading']}")
        lines.append(section["content"].replace("\n", " ").strip())
    csv_path.write_text("\n".join(lines), encoding="utf-8")
    print(f"[FILE] Created CSV: {csv_path} ({csv_path.stat().st_size} bytes)")
    return csv_path


def create_docx_file(temp_dir: Path) -> Path:
    """Create sample DOCX file."""
    try:
        from docx import Document
    except ImportError:
        print("[SKIP] python-docx not installed")
        return None

    doc = Document()
    doc.add_heading(TEST_CONTENT["title"], 0)
    for section in TEST_CONTENT["sections"]:
        doc.add_heading(section["heading"], level=1)
        doc.add_paragraph(section["content"])
    doc.add_heading("Technology Comparison", level=1)
    table = doc.add_table(rows=len(TEST_CONTENT["table_data"]), cols=3)
    table.style = "Table Grid"
    for i, row_data in enumerate(TEST_CONTENT["table_data"]):
        row = table.rows[i]
        for j, cell_text in enumerate(row_data):
            row.cells[j].text = cell_text

    docx_path = temp_dir / "test_document.docx"
    doc.save(str(docx_path))
    print(f"[FILE] Created DOCX: {docx_path} ({docx_path.stat().st_size} bytes)")
    return docx_path


def create_xlsx_file(temp_dir: Path) -> Path:
    """Create sample XLSX file."""
    try:
        from openpyxl import Workbook
    except ImportError:
        print("[SKIP] openpyxl not installed")
        return None

    wb = Workbook()
    ws1 = wb.active
    ws1.title = "Technology Data"
    for row_idx, row_data in enumerate(TEST_CONTENT["table_data"], 1):
        for col_idx, cell_value in enumerate(row_data, 1):
            ws1.cell(row=row_idx, column=col_idx, value=cell_value)

    ws2 = wb.create_sheet("Document Content")
    row = 1
    ws2.cell(row=row, column=1, value=TEST_CONTENT["title"])
    row += 2
    for section in TEST_CONTENT["sections"]:
        ws2.cell(row=row, column=1, value=section["heading"])
        row += 1
        ws2.cell(row=row, column=1, value=section["content"])
        row += 2

    xlsx_path = temp_dir / "test_document.xlsx"
    wb.save(str(xlsx_path))
    print(f"[FILE] Created XLSX: {xlsx_path} ({xlsx_path.stat().st_size} bytes)")
    return xlsx_path


def create_pptx_file(temp_dir: Path) -> Path:
    """Create sample PPTX file."""
    try:
        from pptx import Presentation
    except ImportError:
        print("[SKIP] python-pptx not installed")
        return None

    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = TEST_CONTENT["title"]
    slide.placeholders[1].text = "Multi-Format Ingestion Test"

    content_layout = prs.slide_layouts[1]
    for section in TEST_CONTENT["sections"]:
        slide = prs.slides.add_slide(content_layout)
        slide.shapes.title.text = section["heading"]
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.text = section["content"][:500]

    pptx_path = temp_dir / "test_document.pptx"
    prs.save(str(pptx_path))
    print(f"[FILE] Created PPTX: {pptx_path} ({pptx_path.stat().st_size} bytes)")
    return pptx_path


def create_pdf_file(temp_dir: Path) -> Path:
    """Create sample PDF file."""
    try:
        from reportlab.lib.pagesizes import letter
        from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
        from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle
        from reportlab.lib import colors
    except ImportError:
        print("[SKIP] reportlab not installed")
        return None

    pdf_path = temp_dir / "test_document.pdf"
    doc = SimpleDocTemplate(str(pdf_path), pagesize=letter)
    styles = getSampleStyleSheet()

    title_style = ParagraphStyle("CustomTitle", parent=styles["Heading1"], fontSize=18, spaceAfter=20)
    heading_style = ParagraphStyle("CustomHeading", parent=styles["Heading2"], fontSize=14, spaceAfter=10, spaceBefore=15)

    story = []
    story.append(Paragraph(TEST_CONTENT["title"], title_style))
    story.append(Spacer(1, 12))

    for section in TEST_CONTENT["sections"]:
        story.append(Paragraph(section["heading"], heading_style))
        story.append(Paragraph(section["content"], styles["Normal"]))
        story.append(Spacer(1, 12))

    story.append(Paragraph("Technology Comparison", heading_style))
    table = Table(TEST_CONTENT["table_data"])
    table.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), colors.grey),
        ("TEXTCOLOR", (0, 0), (-1, 0), colors.whitesmoke),
        ("ALIGN", (0, 0), (-1, -1), "CENTER"),
        ("FONTSIZE", (0, 0), (-1, 0), 12),
        ("BOTTOMPADDING", (0, 0), (-1, 0), 12),
        ("GRID", (0, 0), (-1, -1), 1, colors.black),
    ]))
    story.append(table)
    doc.build(story)
    print(f"[FILE] Created PDF: {pdf_path} ({pdf_path.stat().st_size} bytes)")
    return pdf_path


def upload_document(file_path: Path, namespace: str, token: str) -> dict[str, Any]:
    """Upload document and measure timing."""
    headers = {"Authorization": f"Bearer {token}"}

    print(f"\n[UPLOAD] Starting upload: {file_path.name}")
    print(f"[UPLOAD] Namespace: {namespace}")
    print(f"[UPLOAD] Size: {file_path.stat().st_size} bytes")

    start_time = time.time()

    with open(file_path, "rb") as f:
        response = requests.post(
            f"{API_BASE_URL}/api/v1/retrieval/upload",
            headers=headers,
            files={"file": (file_path.name, f)},
            data={"namespace": namespace},
            timeout=900,  # 15 minute timeout
        )

    duration = time.time() - start_time

    result = {
        "format": file_path.suffix,
        "filename": file_path.name,
        "file_size": file_path.stat().st_size,
        "status_code": response.status_code,
        "duration_seconds": round(duration, 2),
        "timestamp": datetime.now().isoformat(),
    }

    if response.status_code == 200:
        data = response.json()
        result.update({
            "success": True,
            "chunks_created": data.get("chunks_created", 0),
            "neo4j_entities": data.get("neo4j_entities", 0),
            "neo4j_relationships": data.get("neo4j_relationships", 0),
            "embeddings_generated": data.get("embeddings_generated", 0),
            "status": data.get("status", "unknown"),
        })
        print(f"[UPLOAD] SUCCESS in {duration:.2f}s")
        print(f"[UPLOAD]   Chunks: {result['chunks_created']}")
        print(f"[UPLOAD]   Entities: {result['neo4j_entities']}")
        print(f"[UPLOAD]   Relations: {result['neo4j_relationships']}")
    else:
        result.update({
            "success": False,
            "error": response.text[:500],
        })
        print(f"[UPLOAD] FAILED ({response.status_code}): {response.text[:200]}")

    return result


def run_tests(formats: list[str], output_file: str = None):
    """Run tests for specified formats."""
    print("=" * 60)
    print("MULTIFORMAT INGESTION TEST")
    print(f"Started: {datetime.now().isoformat()}")
    print("=" * 60)

    # Get auth token
    token = get_auth_token()

    # Create temp directory
    with tempfile.TemporaryDirectory(prefix="aegis_format_test_") as temp_dir:
        temp_path = Path(temp_dir)
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")

        # File creators
        creators = {
            "csv": create_csv_file,
            "docx": create_docx_file,
            "xlsx": create_xlsx_file,
            "pptx": create_pptx_file,
            "pdf": create_pdf_file,
        }

        results = []

        for fmt in formats:
            if fmt not in creators:
                print(f"[SKIP] Unknown format: {fmt}")
                continue

            print(f"\n{'='*60}")
            print(f"TESTING FORMAT: {fmt.upper()}")
            print(f"{'='*60}")

            # Create file
            file_path = creators[fmt](temp_path)
            if file_path is None:
                results.append({
                    "format": fmt,
                    "success": False,
                    "error": "File creation failed (missing dependency)",
                })
                continue

            # Upload
            namespace = f"multiformat_test_{timestamp}_{fmt}"
            result = upload_document(file_path, namespace, token)
            results.append(result)

        # Summary
        print("\n" + "=" * 60)
        print("TEST RESULTS SUMMARY")
        print("=" * 60)
        print(f"{'Format':<10} {'Status':<10} {'Duration':<12} {'Chunks':<8} {'Entities':<10} {'Relations':<10}")
        print("-" * 70)

        for r in results:
            status = "OK" if r.get("success") else "FAILED"
            duration = f"{r.get('duration_seconds', 0):.1f}s"
            chunks = str(r.get("chunks_created", "-"))
            entities = str(r.get("neo4j_entities", "-"))
            relations = str(r.get("neo4j_relationships", "-"))
            print(f"{r['format']:<10} {status:<10} {duration:<12} {chunks:<8} {entities:<10} {relations:<10}")

        # Save results
        if output_file:
            output_path = Path(output_file)
            output_path.parent.mkdir(parents=True, exist_ok=True)
            with open(output_path, "w") as f:
                json.dump({
                    "timestamp": datetime.now().isoformat(),
                    "results": results,
                }, f, indent=2)
            print(f"\n[SAVE] Results saved to: {output_path}")

        return results


def main():
    parser = argparse.ArgumentParser(description="Multi-Format Ingestion Test")
    parser.add_argument("--format", "-f", choices=["csv", "docx", "xlsx", "pptx", "pdf"],
                        help="Test single format")
    parser.add_argument("--all", "-a", action="store_true", help="Test all formats")
    parser.add_argument("--output", "-o", default="data/test_results/multiformat_test.json",
                        help="Output file for results")

    args = parser.parse_args()

    if args.format:
        formats = [args.format]
    elif args.all:
        formats = ["csv", "docx", "xlsx", "pptx", "pdf"]
    else:
        formats = ["csv"]  # Default to quickest test

    run_tests(formats, args.output)


if __name__ == "__main__":
    main()
