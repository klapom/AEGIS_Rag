"""
Test fixtures for AegisRAG integration tests.

Sprint 25 - Feature 25.2: LangGraph Integration Tests
Provides sample documents and common test utilities.
"""

from pathlib import Path

import pytest
from PIL import Image


@pytest.fixture(scope="session")
def fixtures_dir() -> Path:
    """Return the fixtures directory path."""
    return Path(__file__).parent


@pytest.fixture(scope="session")
def sample_text() -> str:
    """Return sample text for testing."""
    return """# Sample Document

This is a sample document for testing the AegisRAG ingestion pipeline.

## Section 1: Introduction
This document contains various elements including text, headings, and structured content.

## Section 2: Technical Details
- Item 1: First bullet point
- Item 2: Second bullet point
- Item 3: Third bullet point

## Section 3: Code Example
```python
def hello_world():
    print("Hello, World!")
```

## Section 4: Entities
The document mentions several important entities:
- **Apple Inc.** is a technology company
- **Microsoft** develops software
- **Google** provides search services

## Section 5: Relationships
Apple Inc. competes with Microsoft.
Microsoft partners with OpenAI.
Google acquired DeepMind.

## Conclusion
This concludes the sample document.
"""


@pytest.fixture(scope="session")
def sample_image_with_text(fixtures_dir: Path) -> Path:
    """
    Create a sample image with text for VLM testing.

    Returns path to generated image file.
    """
    # Create simple image with text
    from PIL import ImageDraw

    img = Image.new("RGB", (800, 400), color="white")
    draw = ImageDraw.Draw(img)

    # Draw some text
    text = "Sample Image\nContains: Blue Rectangle\nRed Circle"
    draw.text((50, 50), text, fill="black")

    # Draw a blue rectangle
    draw.rectangle([50, 150, 250, 250], outline="blue", width=3)

    # Draw a red circle
    draw.ellipse([300, 150, 500, 350], outline="red", width=3)

    # Save image
    image_path = fixtures_dir / "sample_image.png"
    img.save(image_path)

    return image_path


@pytest.fixture(scope="session")
def sample_pdf_path(fixtures_dir: Path, sample_text: str) -> Path:
    """
    Create a sample PDF file for testing.

    Returns path to generated PDF file.

    Note: Requires reportlab (optional dependency).
    If not available, returns None and tests should skip.
    """
    try:
        from reportlab.lib.pagesizes import letter
        from reportlab.lib.units import inch
        from reportlab.pdfgen import canvas
    except ImportError:
        # ReportLab not installed - return None
        return None

    pdf_path = fixtures_dir / "sample_document.pdf"

    # Create simple PDF
    c = canvas.Canvas(str(pdf_path), pagesize=letter)

    # Page 1
    c.setFont("Helvetica-Bold", 16)
    c.drawString(inch, 10.5 * inch, "Sample Document - Page 1")

    c.setFont("Helvetica", 12)
    text_lines = sample_text.split("\n")
    y = 10 * inch
    for line in text_lines[:20]:  # First 20 lines on page 1
        c.drawString(inch, y, line[:80])  # Truncate long lines
        y -= 0.2 * inch
        if y < inch:
            break

    c.showPage()

    # Page 2
    c.setFont("Helvetica-Bold", 16)
    c.drawString(inch, 10.5 * inch, "Sample Document - Page 2")

    c.setFont("Helvetica", 12)
    y = 10 * inch
    for line in text_lines[20:]:  # Remaining lines on page 2
        c.drawString(inch, y, line[:80])
        y -= 0.2 * inch
        if y < inch:
            break

    c.save()

    return pdf_path


@pytest.fixture(scope="session")
def sample_docx_path(fixtures_dir: Path, sample_text: str) -> Path:
    """
    Create a sample DOCX file for testing.

    Returns path to generated DOCX file.

    Note: Requires python-docx (optional dependency).
    If not available, returns None and tests should skip.
    """
    try:
        from docx import Document
    except ImportError:
        # python-docx not installed - return None
        return None

    docx_path = fixtures_dir / "sample_document.docx"

    # Create simple DOCX
    doc = Document()

    # Add title
    doc.add_heading("Sample Document", 0)

    # Add paragraphs from sample text
    for line in sample_text.split("\n"):
        if line.startswith("# "):
            doc.add_heading(line[2:], level=1)
        elif line.startswith("## "):
            doc.add_heading(line[3:], level=2)
        elif line.strip():
            doc.add_paragraph(line)

    doc.save(str(docx_path))

    return docx_path


@pytest.fixture(scope="session")
def sample_pptx_path(fixtures_dir: Path) -> Path:
    """
    Create a sample PPTX file for testing.

    Returns path to generated PPTX file.

    Note: Requires python-pptx (optional dependency).
    If not available, returns None and tests should skip.
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        # python-pptx not installed - return None
        return None

    pptx_path = fixtures_dir / "sample_presentation.pptx"

    # Create simple PPTX
    prs = Presentation()

    # Slide 1: Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Sample Presentation"
    subtitle.text = "Test Document for AegisRAG"

    # Slide 2: Content Slide
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Key Points"
    tf = body_shape.text_frame
    tf.text = "First bullet point"
    p = tf.add_paragraph()
    p.text = "Second bullet point"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Third bullet point"
    p.level = 0

    # Slide 3: Another content slide
    slide = prs.slides.add_slide(bullet_slide_layout)
    title_shape = slide.shapes.title
    body_shape = slide.placeholders[1]
    title_shape.text = "Technical Details"
    tf = body_shape.text_frame
    tf.text = "Contains entities and relationships"

    # Slide 4: Conclusion
    slide = prs.slides.add_slide(bullet_slide_layout)
    title_shape = slide.shapes.title
    body_shape = slide.placeholders[1]
    title_shape.text = "Conclusion"
    tf = body_shape.text_frame
    tf.text = "End of sample presentation"

    # Slide 5: Thank You
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Thank You"
    subtitle.text = "Questions?"

    prs.save(str(pptx_path))

    return pptx_path


@pytest.fixture
def mock_qdrant_client():
    """Mock Qdrant client for testing."""
    from unittest.mock import AsyncMock, MagicMock

    mock_client = AsyncMock()
    mock_client.upsert = AsyncMock(return_value={"status": "acknowledged"})
    mock_client.get_collection = MagicMock(return_value={"points_count": 0})

    return mock_client


@pytest.fixture
def mock_neo4j_driver():
    """Mock Neo4j driver for testing."""
    from unittest.mock import MagicMock

    mock_driver = MagicMock()
    mock_session = MagicMock()
    mock_driver.session.return_value.__enter__.return_value = mock_session
    mock_session.run.return_value = []

    return mock_driver


@pytest.fixture
def mock_llm_response():
    """Mock LLM response for testing."""
    from unittest.mock import AsyncMock

    mock_response = AsyncMock()
    mock_response.content = "Extracted entities: Apple, Microsoft, Google"
    mock_response.provider = "local_ollama"
    mock_response.model = "test-model"
    mock_response.tokens_used = 150
    mock_response.cost_usd = 0.0

    return mock_response
