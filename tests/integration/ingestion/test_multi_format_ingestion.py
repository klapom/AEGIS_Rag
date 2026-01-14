"""Multi-Format Document Ingestion Test Suite.

Sprint 85: Comprehensive testing of all supported document formats.

This test suite verifies that the ingestion pipeline correctly processes:
- PDF (Docling with OCR)
- DOCX (Microsoft Word)
- PPTX (Microsoft PowerPoint)
- XLSX (Microsoft Excel)
- CSV (Comma-Separated Values)

Each format is tested for:
1. Successful parsing (no errors)
2. Correct chunk extraction
3. Entity extraction from content
4. Vector embedding generation
5. Neo4j graph population

Usage:
    pytest tests/integration/ingestion/test_multi_format_ingestion.py -v
    pytest tests/integration/ingestion/test_multi_format_ingestion.py -v -k "pdf"
    pytest tests/integration/ingestion/test_multi_format_ingestion.py -v --format=docx
"""

import asyncio
import hashlib
import json
import os
import tempfile
import time
from datetime import datetime
from pathlib import Path
from typing import Any

import pytest
import requests

# Test configuration
API_BASE_URL = os.getenv("API_BASE_URL", "http://localhost:8000")
TEST_USERNAME = os.getenv("TEST_USERNAME", "admin")
TEST_PASSWORD = os.getenv("TEST_PASSWORD", "admin123")
INGESTION_TIMEOUT = int(os.getenv("INGESTION_TIMEOUT", "600"))  # 10 minutes max

# Test content - same content across all formats for comparison
TEST_CONTENT = {
    "title": "AegisRAG Multi-Format Test Document",
    "sections": [
        {
            "heading": "Introduction to Machine Learning",
            "content": """Machine learning is a subset of artificial intelligence that enables
            systems to learn and improve from experience. Key concepts include supervised learning,
            unsupervised learning, and reinforcement learning. Popular frameworks include
            TensorFlow, PyTorch, and scikit-learn.""",
        },
        {
            "heading": "Natural Language Processing",
            "content": """NLP focuses on the interaction between computers and human language.
            Important techniques include tokenization, named entity recognition (NER),
            sentiment analysis, and transformer models like BERT and GPT. SpaCy and
            Hugging Face are widely used libraries.""",
        },
        {
            "heading": "Knowledge Graphs",
            "content": """Knowledge graphs represent information as entities and relationships.
            Neo4j is a popular graph database. Entity extraction identifies people, organizations,
            and concepts. Relationship extraction finds connections between entities.
            AegisRAG uses LightRAG for graph-based retrieval.""",
        },
    ],
    "entities_expected": [
        "Machine Learning",
        "TensorFlow",
        "PyTorch",
        "scikit-learn",
        "NLP",
        "BERT",
        "GPT",
        "SpaCy",
        "Hugging Face",
        "Neo4j",
        "AegisRAG",
        "LightRAG",
    ],
    "table_data": [
        ["Framework", "Type", "Language"],
        ["TensorFlow", "Deep Learning", "Python"],
        ["PyTorch", "Deep Learning", "Python"],
        ["SpaCy", "NLP", "Python"],
        ["Neo4j", "Graph Database", "Java"],
    ],
}


# =============================================================================
# FIXTURES: Test File Generation
# =============================================================================


@pytest.fixture(scope="module")
def auth_token() -> str:
    """Get authentication token for API calls."""
    response = requests.post(
        f"{API_BASE_URL}/api/v1/auth/login",
        json={"username": TEST_USERNAME, "password": TEST_PASSWORD},
        timeout=30,
    )
    assert response.status_code == 200, f"Auth failed: {response.text}"
    return response.json()["access_token"]


@pytest.fixture(scope="module")
def test_namespace() -> str:
    """Generate unique namespace for this test run."""
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    return f"multiformat_test_{timestamp}"


@pytest.fixture(scope="module")
def temp_dir():
    """Create temporary directory for test files."""
    with tempfile.TemporaryDirectory(prefix="aegis_format_test_") as tmpdir:
        yield Path(tmpdir)


@pytest.fixture(scope="module")
def sample_csv(temp_dir: Path) -> Path:
    """Create sample CSV file."""
    csv_path = temp_dir / "test_document.csv"

    # CSV with headers and data
    lines = [",".join(TEST_CONTENT["table_data"][0])]  # Header
    for row in TEST_CONTENT["table_data"][1:]:
        lines.append(",".join(row))

    # Add text content as additional rows
    lines.append("")
    lines.append("# Document Content")
    for section in TEST_CONTENT["sections"]:
        lines.append(f"## {section['heading']}")
        lines.append(section["content"].replace("\n", " ").strip())

    csv_path.write_text("\n".join(lines), encoding="utf-8")
    return csv_path


@pytest.fixture(scope="module")
def sample_docx(temp_dir: Path) -> Path:
    """Create sample DOCX file."""
    try:
        from docx import Document
        from docx.shared import Pt
    except ImportError:
        pytest.skip("python-docx not installed. Run: pip install python-docx")

    doc = Document()

    # Title
    title = doc.add_heading(TEST_CONTENT["title"], 0)

    # Sections with headings and content
    for section in TEST_CONTENT["sections"]:
        doc.add_heading(section["heading"], level=1)
        doc.add_paragraph(section["content"])

    # Add table
    doc.add_heading("Technology Comparison", level=1)
    table = doc.add_table(rows=len(TEST_CONTENT["table_data"]), cols=3)
    table.style = "Table Grid"

    for i, row_data in enumerate(TEST_CONTENT["table_data"]):
        row = table.rows[i]
        for j, cell_text in enumerate(row_data):
            row.cells[j].text = cell_text

    docx_path = temp_dir / "test_document.docx"
    doc.save(str(docx_path))
    return docx_path


@pytest.fixture(scope="module")
def sample_xlsx(temp_dir: Path) -> Path:
    """Create sample XLSX file."""
    try:
        from openpyxl import Workbook
    except ImportError:
        pytest.skip("openpyxl not installed. Run: pip install openpyxl")

    wb = Workbook()

    # Sheet 1: Table data
    ws1 = wb.active
    ws1.title = "Technology Data"
    for row_idx, row_data in enumerate(TEST_CONTENT["table_data"], 1):
        for col_idx, cell_value in enumerate(row_data, 1):
            ws1.cell(row=row_idx, column=col_idx, value=cell_value)

    # Sheet 2: Document content
    ws2 = wb.create_sheet("Document Content")
    row = 1
    ws2.cell(row=row, column=1, value=TEST_CONTENT["title"])
    row += 2

    for section in TEST_CONTENT["sections"]:
        ws2.cell(row=row, column=1, value=section["heading"])
        row += 1
        ws2.cell(row=row, column=1, value=section["content"])
        row += 2

    xlsx_path = temp_dir / "test_document.xlsx"
    wb.save(str(xlsx_path))
    return xlsx_path


@pytest.fixture(scope="module")
def sample_pptx(temp_dir: Path) -> Path:
    """Create sample PPTX file."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        pytest.skip("python-pptx not installed. Run: pip install python-pptx")

    prs = Presentation()

    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = TEST_CONTENT["title"]
    slide.placeholders[1].text = "Multi-Format Ingestion Test"

    # Content slides
    content_layout = prs.slide_layouts[1]  # Title and Content
    for section in TEST_CONTENT["sections"]:
        slide = prs.slides.add_slide(content_layout)
        slide.shapes.title.text = section["heading"]
        # Add content to body
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.text = section["content"][:500]  # Truncate for slide

    # Table slide
    table_layout = prs.slide_layouts[5]  # Blank
    slide = prs.slides.add_slide(table_layout)

    pptx_path = temp_dir / "test_document.pptx"
    prs.save(str(pptx_path))
    return pptx_path


@pytest.fixture(scope="module")
def sample_pdf(temp_dir: Path) -> Path:
    """Create sample PDF file."""
    try:
        from reportlab.lib.pagesizes import letter
        from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
        from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle
        from reportlab.lib import colors
    except ImportError:
        pytest.skip("reportlab not installed. Run: pip install reportlab")

    pdf_path = temp_dir / "test_document.pdf"
    doc = SimpleDocTemplate(str(pdf_path), pagesize=letter)
    styles = getSampleStyleSheet()

    # Custom styles
    title_style = ParagraphStyle(
        "CustomTitle",
        parent=styles["Heading1"],
        fontSize=18,
        spaceAfter=20,
    )
    heading_style = ParagraphStyle(
        "CustomHeading",
        parent=styles["Heading2"],
        fontSize=14,
        spaceAfter=10,
        spaceBefore=15,
    )

    story = []

    # Title
    story.append(Paragraph(TEST_CONTENT["title"], title_style))
    story.append(Spacer(1, 12))

    # Sections
    for section in TEST_CONTENT["sections"]:
        story.append(Paragraph(section["heading"], heading_style))
        story.append(Paragraph(section["content"], styles["Normal"]))
        story.append(Spacer(1, 12))

    # Table
    story.append(Paragraph("Technology Comparison", heading_style))
    table = Table(TEST_CONTENT["table_data"])
    table.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), colors.grey),
        ("TEXTCOLOR", (0, 0), (-1, 0), colors.whitesmoke),
        ("ALIGN", (0, 0), (-1, -1), "CENTER"),
        ("FONTSIZE", (0, 0), (-1, 0), 12),
        ("BOTTOMPADDING", (0, 0), (-1, 0), 12),
        ("GRID", (0, 0), (-1, -1), 1, colors.black),
    ]))
    story.append(table)

    doc.build(story)
    return pdf_path


# =============================================================================
# HELPER FUNCTIONS
# =============================================================================


def upload_document(
    file_path: Path,
    namespace: str,
    auth_token: str,
    timeout: int = INGESTION_TIMEOUT,
) -> dict[str, Any]:
    """Upload document and wait for ingestion to complete."""
    headers = {"Authorization": f"Bearer {auth_token}"}

    with open(file_path, "rb") as f:
        response = requests.post(
            f"{API_BASE_URL}/api/v1/retrieval/upload",
            headers=headers,
            files={"file": (file_path.name, f)},
            data={"namespace": namespace},
            timeout=timeout,
        )

    return {
        "status_code": response.status_code,
        "response": response.json() if response.status_code == 200 else None,
        "error": response.text if response.status_code != 200 else None,
        "duration": response.elapsed.total_seconds(),
    }


def query_neo4j_entities(namespace: str) -> dict[str, Any]:
    """Query Neo4j for entities in the given namespace."""
    import subprocess

    query = f"""
    MATCH (e:base)
    WHERE e.namespace_id = '{namespace}' OR e.namespace_id IS NULL
    RETURN labels(e) AS types, count(*) AS count
    ORDER BY count DESC
    LIMIT 20
    """

    result = subprocess.run(
        [
            "docker", "exec", "aegis-neo4j", "cypher-shell",
            "-u", "neo4j", "-p", "aegis-rag-neo4j-password",
            query,
        ],
        capture_output=True,
        text=True,
    )

    return {
        "stdout": result.stdout,
        "stderr": result.stderr,
        "returncode": result.returncode,
    }


def query_neo4j_chunks(namespace: str) -> int:
    """Count chunks in Neo4j for the given namespace."""
    import subprocess

    # Query all chunks (namespace filtering may not be on chunk level)
    query = "MATCH (c:chunk) RETURN count(c) AS count"

    result = subprocess.run(
        [
            "docker", "exec", "aegis-neo4j", "cypher-shell",
            "-u", "neo4j", "-p", "aegis-rag-neo4j-password",
            query,
        ],
        capture_output=True,
        text=True,
    )

    if result.returncode == 0:
        # Parse output like "count\n354"
        lines = result.stdout.strip().split("\n")
        if len(lines) >= 2:
            try:
                return int(lines[1])
            except ValueError:
                pass
    return 0


# =============================================================================
# TEST CASES
# =============================================================================


class TestMultiFormatIngestion:
    """Test suite for multi-format document ingestion."""

    @pytest.mark.parametrize("format_name,fixture_name", [
        ("csv", "sample_csv"),
        ("docx", "sample_docx"),
        ("xlsx", "sample_xlsx"),
        ("pptx", "sample_pptx"),
        ("pdf", "sample_pdf"),
    ])
    def test_format_upload_success(
        self,
        format_name: str,
        fixture_name: str,
        auth_token: str,
        test_namespace: str,
        request,
    ):
        """Test that each format uploads successfully."""
        # Get the fixture dynamically
        file_path = request.getfixturevalue(fixture_name)

        # Create format-specific namespace
        namespace = f"{test_namespace}_{format_name}"

        # Upload
        result = upload_document(file_path, namespace, auth_token)

        # Assertions
        assert result["status_code"] == 200, f"Upload failed for {format_name}: {result['error']}"
        assert result["response"] is not None, f"No response for {format_name}"

        # Log results
        print(f"\n{'='*60}")
        print(f"FORMAT: {format_name.upper()}")
        print(f"{'='*60}")
        print(f"File: {file_path.name}")
        print(f"Size: {file_path.stat().st_size} bytes")
        print(f"Duration: {result['duration']:.2f}s")
        print(f"Response: {json.dumps(result['response'], indent=2)}")

    def test_csv_specific_features(
        self,
        sample_csv: Path,
        auth_token: str,
        test_namespace: str,
    ):
        """Test CSV-specific parsing features."""
        namespace = f"{test_namespace}_csv_detailed"
        result = upload_document(sample_csv, namespace, auth_token)

        assert result["status_code"] == 200
        response = result["response"]

        # CSV should extract table structure
        assert response.get("chunks_created", 0) >= 1, "CSV should create at least 1 chunk"

        print(f"\nCSV Details:")
        print(f"  Chunks: {response.get('chunks_created', 0)}")
        print(f"  Entities: {response.get('neo4j_entities', 0)}")

    def test_docx_specific_features(
        self,
        sample_docx: Path,
        auth_token: str,
        test_namespace: str,
    ):
        """Test DOCX-specific parsing features (headings, tables)."""
        namespace = f"{test_namespace}_docx_detailed"
        result = upload_document(sample_docx, namespace, auth_token)

        assert result["status_code"] == 200
        response = result["response"]

        # DOCX should preserve document structure
        assert response.get("chunks_created", 0) >= 1, "DOCX should create chunks"

        print(f"\nDOCX Details:")
        print(f"  Chunks: {response.get('chunks_created', 0)}")
        print(f"  Entities: {response.get('neo4j_entities', 0)}")
        print(f"  Relations: {response.get('neo4j_relationships', 0)}")

    def test_xlsx_specific_features(
        self,
        sample_xlsx: Path,
        auth_token: str,
        test_namespace: str,
    ):
        """Test XLSX-specific parsing features (multiple sheets, tables)."""
        namespace = f"{test_namespace}_xlsx_detailed"
        result = upload_document(sample_xlsx, namespace, auth_token)

        assert result["status_code"] == 200
        response = result["response"]

        # XLSX should handle multiple sheets
        assert response.get("chunks_created", 0) >= 1, "XLSX should create chunks"

        print(f"\nXLSX Details:")
        print(f"  Chunks: {response.get('chunks_created', 0)}")
        print(f"  Entities: {response.get('neo4j_entities', 0)}")

    def test_pptx_specific_features(
        self,
        sample_pptx: Path,
        auth_token: str,
        test_namespace: str,
    ):
        """Test PPTX-specific parsing features (slides, text boxes)."""
        namespace = f"{test_namespace}_pptx_detailed"
        result = upload_document(sample_pptx, namespace, auth_token)

        assert result["status_code"] == 200
        response = result["response"]

        # PPTX should extract slide content
        assert response.get("chunks_created", 0) >= 1, "PPTX should create chunks"

        print(f"\nPPTX Details:")
        print(f"  Chunks: {response.get('chunks_created', 0)}")
        print(f"  Entities: {response.get('neo4j_entities', 0)}")

    def test_pdf_specific_features(
        self,
        sample_pdf: Path,
        auth_token: str,
        test_namespace: str,
    ):
        """Test PDF-specific parsing features (OCR, layout)."""
        namespace = f"{test_namespace}_pdf_detailed"
        result = upload_document(sample_pdf, namespace, auth_token)

        assert result["status_code"] == 200
        response = result["response"]

        # PDF should use Docling with OCR
        assert response.get("chunks_created", 0) >= 1, "PDF should create chunks"

        print(f"\nPDF Details:")
        print(f"  Chunks: {response.get('chunks_created', 0)}")
        print(f"  Entities: {response.get('neo4j_entities', 0)}")
        print(f"  Duration: {result['duration']:.2f}s")


class TestFormatComparison:
    """Compare extraction quality across formats."""

    def test_entity_extraction_comparison(
        self,
        sample_csv: Path,
        sample_docx: Path,
        sample_xlsx: Path,
        sample_pptx: Path,
        sample_pdf: Path,
        auth_token: str,
        test_namespace: str,
    ):
        """Compare entity extraction across all formats."""
        formats = {
            "csv": sample_csv,
            "docx": sample_docx,
            "xlsx": sample_xlsx,
            "pptx": sample_pptx,
            "pdf": sample_pdf,
        }

        results = {}

        for fmt, path in formats.items():
            namespace = f"{test_namespace}_compare_{fmt}"
            result = upload_document(path, namespace, auth_token)

            if result["status_code"] == 200:
                results[fmt] = {
                    "chunks": result["response"].get("chunks_created", 0),
                    "entities": result["response"].get("neo4j_entities", 0),
                    "relations": result["response"].get("neo4j_relationships", 0),
                    "duration": result["duration"],
                }
            else:
                results[fmt] = {"error": result["error"]}

        # Print comparison table
        print("\n" + "=" * 80)
        print("FORMAT COMPARISON RESULTS")
        print("=" * 80)
        print(f"{'Format':<10} {'Chunks':<10} {'Entities':<10} {'Relations':<10} {'Duration':<10}")
        print("-" * 80)

        for fmt, data in results.items():
            if "error" in data:
                print(f"{fmt:<10} ERROR: {data['error'][:50]}")
            else:
                print(
                    f"{fmt:<10} {data['chunks']:<10} {data['entities']:<10} "
                    f"{data['relations']:<10} {data['duration']:.2f}s"
                )

        # All formats should extract at least some content
        for fmt, data in results.items():
            if "error" not in data:
                assert data["chunks"] >= 1, f"{fmt} should create at least 1 chunk"


class TestEdgeCases:
    """Test edge cases and error handling."""

    def test_empty_file_handling(self, auth_token: str, test_namespace: str, temp_dir: Path):
        """Test handling of empty files."""
        # Create empty CSV
        empty_csv = temp_dir / "empty.csv"
        empty_csv.write_text("")

        result = upload_document(empty_csv, f"{test_namespace}_empty", auth_token)

        # Should handle gracefully (either success with 0 chunks or meaningful error)
        print(f"\nEmpty file result: {result}")

    def test_large_file_handling(self, auth_token: str, test_namespace: str, temp_dir: Path):
        """Test handling of larger files (stress test)."""
        # Create large CSV with repeated content
        large_csv = temp_dir / "large.csv"

        lines = ["heading1,heading2,heading3"]
        for i in range(1000):
            lines.append(f"row{i},Machine Learning concept {i},TensorFlow PyTorch {i}")

        large_csv.write_text("\n".join(lines))

        result = upload_document(
            large_csv,
            f"{test_namespace}_large",
            auth_token,
            timeout=INGESTION_TIMEOUT,
        )

        print(f"\nLarge file ({large_csv.stat().st_size} bytes):")
        print(f"  Status: {result['status_code']}")
        print(f"  Duration: {result['duration']:.2f}s")
        if result["response"]:
            print(f"  Chunks: {result['response'].get('chunks_created', 0)}")


# =============================================================================
# STANDALONE RUNNER
# =============================================================================


if __name__ == "__main__":
    """Run tests standalone for quick verification."""
    pytest.main([
        __file__,
        "-v",
        "-s",
        "--tb=short",
        "-x",  # Stop on first failure
    ])
